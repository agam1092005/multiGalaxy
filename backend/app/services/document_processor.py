"""
Document processing service for Multi-Galaxy-Note.
Handles file upload, text extraction, OCR, and document indexing.
"""

import os
import io
import uuid
from typing import List, Dict, Any, Optional, BinaryIO
from pathlib import Path
import asyncio
import logging

import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
import pytesseract
from PIL import Image
from pdf2image import convert_from_bytes

from ..core.config import get_settings
from .rag_system import RAGSystem

logger = logging.getLogger(__name__)
settings = get_settings()

class DocumentProcessor:
    """Service for processing uploaded documents and extracting text content."""
    
    def __init__(self):
        self.rag_system = RAGSystem()
        
        # Supported file types
        self.supported_types = {
            'application/pdf': self._extract_pdf_text,
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': self._extract_docx_text,
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': self._extract_pptx_text,
            'image/jpeg': self._extract_image_text,
            'image/png': self._extract_image_text,
            'image/tiff': self._extract_image_text,
            'text/plain': self._extract_text_file
        }
    
    async def process_document(
        self, 
        file_content: bytes, 
        filename: str, 
        content_type: str,
        user_id: str
    ) -> Dict[str, Any]:
        """
        Process uploaded document and extract text content.
        
        Args:
            file_content: Binary content of the file
            filename: Original filename
            content_type: MIME type of the file
            user_id: ID of the user uploading the document
            
        Returns:
            Dictionary containing processed document information
        """
        try:
            # Validate file type
            if content_type not in self.supported_types:
                raise ValueError(f"Unsupported file type: {content_type}")
            
            # Generate unique document ID
            doc_id = str(uuid.uuid4())
            
            # Extract text content
            extractor = self.supported_types[content_type]
            text_content = await extractor(file_content)
            
            if not text_content.strip():
                raise ValueError("No text content could be extracted from the document")
            
            # Add to RAG system
            rag_metadata = {
                "filename": filename,
                "user_id": user_id,
                "content_type": content_type,
                "upload_date": doc_id  # Using doc_id as timestamp reference
            }
            
            rag_result = await self.rag_system.add_document(
                document_id=doc_id,
                content=text_content,
                metadata=rag_metadata
            )
            
            # Save original file (in production, this would go to cloud storage)
            file_path = await self._save_file(file_content, doc_id, filename)
            
            return {
                "document_id": doc_id,
                "filename": filename,
                "content_type": content_type,
                "text_content": text_content,
                "chunk_count": rag_result.get('chunks_created', 0),
                "subject": rag_result.get('subject', 'general'),
                "file_path": file_path,
                "status": "processed",
                "rag_collections": rag_result.get('collections', [])
            }
            
        except Exception as e:
            logger.error(f"Error processing document {filename}: {str(e)}")
            raise ValueError(f"Failed to process document: {str(e)}")
    
    async def _extract_pdf_text(self, file_content: bytes) -> str:
        """Extract text from PDF file."""
        try:
            # First try direct text extraction
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
            text_content = ""
            
            for page in pdf_reader.pages:
                text_content += page.extract_text() + "\n"
            
            # If no text found, use OCR
            if not text_content.strip():
                text_content = await self._ocr_pdf(file_content)
            
            return text_content
            
        except Exception as e:
            logger.error(f"Error extracting PDF text: {str(e)}")
            # Fallback to OCR
            return await self._ocr_pdf(file_content)
    
    async def _ocr_pdf(self, file_content: bytes) -> str:
        """Extract text from PDF using OCR."""
        try:
            # Convert PDF to images
            images = convert_from_bytes(file_content)
            text_content = ""
            
            for image in images:
                # Perform OCR on each page
                page_text = pytesseract.image_to_string(image)
                text_content += page_text + "\n"
            
            return text_content
            
        except Exception as e:
            logger.error(f"Error performing OCR on PDF: {str(e)}")
            raise ValueError("Failed to extract text from PDF using OCR")
    
    async def _extract_docx_text(self, file_content: bytes) -> str:
        """Extract text from DOCX file."""
        try:
            doc = DocxDocument(io.BytesIO(file_content))
            text_content = ""
            
            for paragraph in doc.paragraphs:
                text_content += paragraph.text + "\n"
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text_content += cell.text + " "
                    text_content += "\n"
            
            return text_content
            
        except Exception as e:
            logger.error(f"Error extracting DOCX text: {str(e)}")
            raise ValueError("Failed to extract text from DOCX file")
    
    async def _extract_pptx_text(self, file_content: bytes) -> str:
        """Extract text from PPTX file."""
        try:
            prs = Presentation(io.BytesIO(file_content))
            text_content = ""
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + "\n"
            
            return text_content
            
        except Exception as e:
            logger.error(f"Error extracting PPTX text: {str(e)}")
            raise ValueError("Failed to extract text from PPTX file")
    
    async def _extract_image_text(self, file_content: bytes) -> str:
        """Extract text from image using OCR."""
        try:
            image = Image.open(io.BytesIO(file_content))
            text_content = pytesseract.image_to_string(image)
            return text_content
            
        except Exception as e:
            logger.error(f"Error extracting text from image: {str(e)}")
            raise ValueError("Failed to extract text from image using OCR")
    
    async def _extract_text_file(self, file_content: bytes) -> str:
        """Extract text from plain text file."""
        try:
            return file_content.decode('utf-8')
        except UnicodeDecodeError:
            # Try other encodings
            for encoding in ['latin-1', 'cp1252']:
                try:
                    return file_content.decode(encoding)
                except UnicodeDecodeError:
                    continue
            raise ValueError("Failed to decode text file")
    

    
    async def _save_file(self, file_content: bytes, doc_id: str, filename: str) -> str:
        """Save file to storage (local for development, cloud for production)."""
        # Create uploads directory if it doesn't exist
        upload_dir = Path("uploads")
        upload_dir.mkdir(exist_ok=True)
        
        # Generate safe filename
        file_extension = Path(filename).suffix
        safe_filename = f"{doc_id}{file_extension}"
        file_path = upload_dir / safe_filename
        
        # Save file
        with open(file_path, 'wb') as f:
            f.write(file_content)
        
        return str(file_path)
    
    async def search_documents(
        self, 
        query: str, 
        user_id: str, 
        subject: Optional[str] = None,
        limit: int = 5
    ) -> List[Dict[str, Any]]:
        """
        Search for relevant document chunks based on query using the RAG system.
        
        Args:
            query: Search query
            user_id: User ID to filter documents
            subject: Optional subject filter
            limit: Maximum number of results
            
        Returns:
            List of relevant document chunks with metadata
        """
        try:
            return await self.rag_system.search(
                query=query,
                user_id=user_id,
                subject=subject,
                limit=limit,
                similarity_threshold=0.6
            )
            
        except Exception as e:
            logger.error(f"Error searching documents: {str(e)}")
            return []
    
    async def get_user_documents(self, user_id: str) -> List[Dict[str, Any]]:
        """Get all documents for a user using the RAG system."""
        try:
            # Use RAG system to get documents
            search_results = await self.rag_system.search(
                query="",  # Empty query to get all documents
                user_id=user_id,
                limit=1000,  # Large limit to get all documents
                similarity_threshold=0.0  # Low threshold to include all
            )
            
            # Group by document_id
            documents = {}
            for result in search_results:
                metadata = result.get('metadata', {})
                doc_id = metadata.get('document_id')
                if doc_id and doc_id not in documents:
                    documents[doc_id] = {
                        "document_id": doc_id,
                        "filename": metadata.get('filename', 'Unknown'),
                        "content_type": metadata.get('content_type', 'Unknown'),
                        "subject": metadata.get('subject', 'general'),
                        "chunk_count": 0
                    }
                if doc_id:
                    documents[doc_id]["chunk_count"] += 1
            
            return list(documents.values())
            
        except Exception as e:
            logger.error(f"Error getting user documents: {str(e)}")
            return []
    
    async def delete_document(self, document_id: str, user_id: str) -> bool:
        """Delete a document and all its chunks using the RAG system."""
        try:
            # Delete from RAG system
            rag_deleted = await self.rag_system.delete_document(document_id, user_id)
            
            # Delete file (in production, delete from cloud storage)
            try:
                # Find and delete the file
                upload_dir = Path("uploads")
                for file_path in upload_dir.glob(f"{document_id}.*"):
                    file_path.unlink()
            except Exception as e:
                logger.warning(f"Could not delete file for document {document_id}: {str(e)}")
            
            return rag_deleted
            
        except Exception as e:
            logger.error(f"Error deleting document {document_id}: {str(e)}")
            return False
    
    async def get_document_context(
        self,
        query: str,
        user_id: str,
        subject: Optional[str] = None,
        max_context_length: int = 4000
    ) -> Dict[str, Any]:
        """
        Get relevant document context for a query using the RAG system.
        
        Args:
            query: The query to find context for
            user_id: User ID for personalized results
            subject: Optional subject filter
            max_context_length: Maximum total length of context
            
        Returns:
            Dictionary with formatted context and metadata
        """
        try:
            return await self.rag_system.get_context_for_query(
                query=query,
                user_id=user_id,
                subject=subject,
                max_context_length=max_context_length
            )
        except Exception as e:
            logger.error(f"Error getting document context: {str(e)}")
            return {
                'context': '',
                'sources': [],
                'total_chunks': 0,
                'subjects_covered': [],
                'error': str(e)
            }